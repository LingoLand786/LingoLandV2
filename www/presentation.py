from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR

# Create presentation
prs = Presentation()

# Define modern color scheme
PRIMARY_COLOR = RGBColor(41, 128, 185)    # Modern blue
SECONDARY_COLOR = RGBColor(52, 152, 219)  # Lighter blue
ACCENT_COLOR = RGBColor(231, 76, 60)      # Red accent
TEXT_COLOR = RGBColor(44, 62, 80)         # Dark blue-gray
BG_COLOR = RGBColor(236, 240, 241)        # Light gray
WHITE = RGBColor(255, 255, 255)

# Slide dimensions
SLIDE_WIDTH = Inches(16)
hbbb
def create_modern_background(slide):
    """Create a modern gradient background."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[0].color.rgb = PRIMARY_COLOR
    fill.gradient_stops[1].position = 1
    fill.gradient_stops[1].color.rgb = SECONDARY_COLOR

def add_decorative_element(slide):
    """Add decorative elements to slides."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(0.5),
        Inches(15),
        Inches(0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()

def format_title(title_shape, title_text):
    """Format title with modern styling."""
    title_shape.text = title_text
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)

def format_content_paragraph(paragraph, text, level=0):
    """Format content with modern styling."""
    paragraph.text = text
    paragraph.font.name = 'Segoe UI'
    paragraph.font.size = Pt(24) if level == 0 else Pt(20)
    paragraph.font.color.rgb = WHITE
    paragraph.font.bold = True if level == 0 else False
    paragraph.level = level
    paragraph.space_after = Pt(12)

def add_slide(title, content_lines, slide_type='content'):
    """Create a modern slide with styling."""
    if slide_type == 'title':
        slide_layout = prs.slide_layouts[0]
    else:
        slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(slide_layout)
    create_modern_background(slide)
    add_decorative_element(slide)
    
    # Title
    title_shape = slide.shapes.title
    format_title(title_shape, title)
    
    # Content
    if content_lines:
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()
        
        for line in content_lines:
            p = tf.add_paragraph()
            format_content_paragraph(p, line)

# Title slide
add_slide(
    "LingoLand — Aplicação para Aprender Português",
    [
        "Sitara Liaqat & Ana Fillipa",
        "Curso Profissional de Programador/a de Informática",
        "Orientadora: Prof. Catarina Dias",
        "Agrupamento de Escolas Amato Lusitano",
        "Data: [Inserir data da apresentação]"
    ],
    'title'
)

# Content slides with enhanced structure
slides_content = [
    ("Agenda", [
        "• Apresentação Pessoal e do Projeto",
        "• Enquadramento e Objetivos",
        "• Desenvolvimento / Metodologia",
        "• Demonstração Prática",
        "• Conclusões e Reflexão",
        "• Dificuldades e Superações",
        "• Perspetivas Futuras",
        "• Agradecimentos",
        "• Espaço para Questões"
    ]),
    ("Apresentação Pessoal e do Projeto", [
        "• Sitara & Ana – Curso Profissional de Informática",
        "• Projeto: LingoLand (site + aplicação móvel)",
        "• Motivação pessoal: dificuldades na aprendizagem do português",
        "• Projeto de Grupo — Prático"
    ]),
    ("Enquadramento e Objetivos", [
        "• Relevância: apoio à integração de imigrantes",
        "• Objetivo geral: ensinar português de forma acessível",
        "• Objetivos específicos: lições com áudio, gamificação, suporte multilíngue",
        "• Público-alvo: iniciantes absolutos, falantes de Urdu, Pashto, etc."
    ]),
    ("Desenvolvimento / Metodologia", [
        "• Planeamento → Desenvolvimento → Testes",
        "• Tecnologias: Cordova, Firebase, gTTS, HTML, CSS, JS",
        "• Ferramentas: VS Code, Android Studio, Netlify",
        "• Trabalho em equipa e aprendizagem ativa"
    ]),
    ("Demonstração Prática", [
        "• Tela inicial e criação de conta",
        "• Lições com áudio, imagens e exercícios",
        "• Funcionalidades: Palavras Diárias, Comunidade, Contato",
        "• Demonstração ao vivo ou com vídeo/screenshot"
    ]),
    ("Conclusões e Reflexão", [
        "• Desenvolvimento técnico e pessoal",
        "• Impacto real e social do projeto",
        "• Valor pessoal da experiência como imigrante",
        "• Projeto como exemplo do percurso formativo"
    ]),
    ("Dificuldades e Superações", [
        "• Plataformas limitadas (Thunkable)",
        "• Desafios técnicos com Cordova e GitHub Pages",
        "• Migração para Netlify",
        "• Ajuda de colegas, professores e tutoriais"
    ]),
    ("Perspetivas Futuras", [
        "• Adicionar mais idiomas e funcionalidades",
        "• Expandir conteúdo (lições, jogos, quizzes)",
        "• Utilidade em centros de apoio e escolas",
        "• Projeto com valor no mercado e no portefólio"
    ]),
    ("Agradecimentos", [
        "• Professores: Cristiana Silva, Catarina Dias, Helder, Fernando",
        "• Colega Pedro (apoio técnico)",
        "• Famílias",
        "• Agrupamento de Escolas Amato Lusitano"
    ]),
    ("Espaço para Questões", [
        "• Obrigada pela atenção!",
        "• Estamos disponíveis para responder a quaisquer questões."
    ])
]

for title, content in slides_content:
    add_slide(title, content)

# Save presentation with enhanced filename
output_path = "LingoLand_Apresentacao_PAP_Profissional_Enhanced.pptx"
prs.save(output_path)

print(f"Enhanced presentation saved to {output_path}")
